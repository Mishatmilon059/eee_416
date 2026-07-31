#!/usr/bin/env python3
"""Structural QA for docs/BraillePresentation.pptx.

    python3 tools/test_deck.py

Why this is not a render-based check: LibreOffice in some sandboxes refuses to
open ANY pptxgenjs output -- a one-textbox deck fails identically -- so a render
QA there proves nothing about the file. These checks read the package directly,
which is both more precise about geometry and immune to that.

Checks the defects that actually reach a viewer:
  * a shape positioned off the slide, or crossing its edge
  * a slide whose background is not pure white (a stated hard requirement)
  * text boxes overlapping each other
  * missing speaker notes
  * diagram slides that lost their editable SVG
  * numbers on the slides disagreeing with models/metrics.json
"""
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "docs" / "BraillePresentation.pptx"
EXPECTED_SLIDES = 15
DIAGRAM_SLIDES = range(2, 14)          # 1-based: slides 2..13
TOL = Emu(9525 * 4)                    # ~4 px of slack on edge checks

failures = []


def check(name, ok, detail=""):
    print(f"  {'PASS' if ok else 'FAIL'}  {name}" + (f"  -- {detail}" if detail and not ok else ""))
    if not ok:
        failures.append(name)
    return ok


def overlaps(a, b):
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def main():
    if not PPTX.exists():
        sys.exit(f"{PPTX.relative_to(ROOT)} not found -- run node tools/build_deck.js")

    prs = Presentation(PPTX)
    SW, SH = prs.slide_width, prs.slide_height

    print(f"deck   : {PPTX.relative_to(ROOT)}  ({PPTX.stat().st_size/1024:.0f} KB)")
    print(f"canvas : {SW/914400:.3f} x {SH/914400:.3f} inches "
          f"(ratio {SW/SH:.4f}, 16:9 = {16/9:.4f})")

    print("\n-- structure --")
    check(f"{EXPECTED_SLIDES} slides", len(prs.slides.__iter__.__self__._sldIdLst) == EXPECTED_SLIDES
          if hasattr(prs.slides, "_sldIdLst") else len(prs.slides._sldIdLst) == EXPECTED_SLIDES,
          f"got {len(prs.slides._sldIdLst)}")
    check("canvas is 16:9", abs(SW / SH - 16 / 9) < 0.01, f"{SW/SH:.4f}")

    print("\n-- geometry: nothing off-slide --")
    off = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            l, t = sh.left, sh.top
            r, b = l + (sh.width or 0), t + (sh.height or 0)
            if l < -TOL or t < -TOL or r > SW + TOL or b > SH + TOL:
                off.append(f"slide{i} '{(sh.name or '?')[:18]}' "
                           f"({l/914400:.2f},{t/914400:.2f})-({r/914400:.2f},{b/914400:.2f})")
    check("every shape is within the slide", not off, "; ".join(off[:3]))

    print("\n-- text boxes do not overlap --")
    clashes = []
    for i, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            if sh.left is None:
                continue
            boxes.append((sh.left, sh.top, sh.left + (sh.width or 0),
                          sh.top + (sh.height or 0), sh.text_frame.text.strip()[:22]))
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                if overlaps(boxes[a], boxes[b]):
                    clashes.append(f"slide{i}: '{boxes[a][4]}' x '{boxes[b][4]}'")
    check("no two text boxes overlap", not clashes, "; ".join(clashes[:3]))

    print("\n-- pure white backgrounds (hard requirement) --")
    with zipfile.ZipFile(PPTX) as z:
        bad_bg, missing_bg = [], []
        for i in range(1, EXPECTED_SLIDES + 1):
            xml = z.read(f"ppt/slides/slide{i}.xml").decode("utf-8")
            m = re.search(r"<p:bg>.*?</p:bg>", xml, re.S)
            if not m:
                missing_bg.append(f"slide{i}")
                continue
            colors = set(re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"/>', m.group(0)))
            if colors - {"FFFFFF"}:
                bad_bg.append(f"slide{i}={','.join(sorted(colors))}")
        check("every slide background is explicitly FFFFFF",
              not bad_bg and not missing_bg,
              f"non-white: {bad_bg}  missing: {missing_bg}")

        print("\n-- diagram slides keep their editable SVG --")
        no_svg = []
        for i in DIAGRAM_SLIDES:
            xml = z.read(f"ppt/slides/slide{i}.xml").decode("utf-8")
            if "svgBlip" not in xml:
                no_svg.append(f"slide{i}")
        check(f"all {len(list(DIAGRAM_SLIDES))} diagram slides carry an svgBlip",
              not no_svg, f"missing on {no_svg}")

        svgs = [n for n in z.namelist() if n.startswith("ppt/media/") and n.endswith(".svg")]
        pngs = [n for n in z.namelist() if n.startswith("ppt/media/") and n.endswith(".png")]
        check("SVG and PNG counts match (fallback intact)",
              len(svgs) == len(pngs) == 12, f"{len(svgs)} svg, {len(pngs)} png")

    # The usual way to catch overflow is to render and look. LibreOffice cannot
    # open pptxgenjs output in this sandbox, so instead the text is measured
    # against its box using Liberation Sans -- metric-compatible with Arial, the
    # font the deck actually specifies -- and wrapping is simulated.
    print("\n-- text fits inside its box (measured, since rendering is unavailable) --")
    try:
        from PIL import ImageFont
        FONT_PATH = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
        FONT_BOLD = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
        EMU_IN = 914400

        def wrapped_height_in(text, pt, width_in, bold):
            f = ImageFont.truetype(FONT_BOLD if bold else FONT_PATH, size=max(8, int(pt * 4)))
            scale = pt / max(8, int(pt * 4))          # measure big, scale down
            max_px = width_in * 72 / scale            # box width in the measured units
            lines = 0
            for para in text.split("\n"):
                words, cur = para.split(), ""
                if not words:
                    lines += 1
                    continue
                n = 0
                for w in words:
                    trial = (cur + " " + w).strip()
                    if f.getlength(trial) <= max_px or not cur:
                        cur = trial
                    else:
                        n += 1
                        cur = w
                lines += n + 1
            return lines * pt * 1.22 / 72             # 1.22 line spacing, pt -> inches

        overflow = []
        for i, slide in enumerate(prs.slides, 1):
            for sh in slide.shapes:
                if not sh.has_text_frame or not sh.text_frame.text.strip():
                    continue
                if sh.width is None or sh.height is None:
                    continue
                para = sh.text_frame.paragraphs[0]
                run = para.runs[0] if para.runs else None
                if run is None or run.font.size is None:
                    continue
                pt = run.font.size.pt
                bold = bool(run.font.bold)
                need = wrapped_height_in(sh.text_frame.text, pt, sh.width / EMU_IN, bold)
                have = sh.height / EMU_IN
                if need > have * 1.15:                # 15% slack
                    overflow.append(f"slide{i} '{sh.text_frame.text.strip()[:24]}' "
                                    f"needs {need:.2f}in, box {have:.2f}in")
        check("no text box is too small for its text", not overflow,
              "; ".join(overflow[:4]))
    except ImportError:
        print("  SKIP  Pillow not available")

    print("\n-- speaker notes --")
    no_notes = [i for i, s in enumerate(prs.slides, 1)
                if not (s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())]
    check("every slide has speaker notes", not no_notes, f"missing on {no_notes}")

    print("\n-- numbers agree with models/metrics.json --")
    M = json.loads((ROOT / "models" / "metrics.json").read_text())
    all_text = "\n".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                         if sh.has_text_frame)
    all_text += "\n".join(s.notes_slide.notes_text_frame.text for s in prs.slides
                          if s.has_notes_slide)
    wants = [
        (f"{M['params']:,}", "parameter count"),
        (f"{M['tflite_bytes']:,}", "quantized size"),
        (f"{M['test']['teaching_combined']*100:.1f}", "teaching accuracy"),
    ]
    for val, label in wants:
        check(f"{label} ({val}) appears in the deck", val in all_text)

    bmap = json.loads((ROOT / "data" / "braille_map.json").read_text())
    n_ver = sum(1 for l in bmap["letters"] if l.get("verified"))
    check(f"verified-letter count ({n_ver}) is stated", f"{n_ver} of 50" in all_text)

    print("\n-- no placeholder text left behind --")
    leftovers = re.findall(r"lorem|ipsum|\bTODO\b|\[insert|XXXX", all_text, re.I)
    check("no placeholder text", not leftovers, str(set(leftovers)))

    print("\n" + "=" * 58)
    if failures:
        print(f"{len(failures)} check(s) FAILED: {', '.join(failures)}")
        return 1
    print("OK - deck passes structural QA")
    return 0


if __name__ == "__main__":
    sys.exit(main())
